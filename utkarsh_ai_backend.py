"""
╔══════════════════════════════════════════════════════════════╗
║              UTKARSH AI — FastAPI Backend                    ║
║  Stack: FastAPI + LangChain + python-pptx + ChromaDB         ║
╚══════════════════════════════════════════════════════════════╝

SETUP:
  pip install fastapi uvicorn python-pptx langchain langchain-openai
              chromadb pypdf python-docx pillow python-multipart
              anthropic openai

RUN:
  uvicorn utkarsh_ai_backend:app --reload --port 8000
"""

import os, io, uuid, json, base64
from pathlib import Path
from typing import Optional

from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from pydantic import BaseModel

# ── PPT ──────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── AI / LangChain ───────────────────────────────────────────────
from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from langchain.chains import RetrievalQA
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
from langchain_community.vectorstores import Chroma
from langchain.schema import HumanMessage, AIMessage, SystemMessage

# ── Image / OCR ──────────────────────────────────────────────────
from PIL import Image
import anthropic   # for vision calls

# ─────────────────────────────────────────────────────────────────
app = FastAPI(title="Utkarsh AI API", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

UPLOAD_DIR = Path("uploads");  UPLOAD_DIR.mkdir(exist_ok=True)
OUTPUT_DIR = Path("outputs");  OUTPUT_DIR.mkdir(exist_ok=True)
CHROMA_DIR = Path("chroma_db"); CHROMA_DIR.mkdir(exist_ok=True)

OPENAI_API_KEY   = os.getenv("OPENAI_API_KEY", "")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY", "")

llm        = ChatOpenAI(model="gpt-4o", api_key=OPENAI_API_KEY, temperature=0.7)
embeddings = OpenAIEmbeddings(api_key=OPENAI_API_KEY)
ant_client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)


# ═══════════════════════════════════════════════════════════════
#  MODELS
# ═══════════════════════════════════════════════════════════════
class ChatRequest(BaseModel):
    message: str
    history: list[dict] = []
    session_id: str = "default"

class PPTRequest(BaseModel):
    topic: str
    num_slides: int = 8
    style: str = "professional"   # professional | creative | minimal


# ═══════════════════════════════════════════════════════════════
#  HEALTH
# ═══════════════════════════════════════════════════════════════
@app.get("/")
def root():
    return {"status": "Utkarsh AI is running 🚀", "version": "1.0.0"}


# ═══════════════════════════════════════════════════════════════
#  1.  CHAT  (streaming)
# ═══════════════════════════════════════════════════════════════
@app.post("/chat")
async def chat(req: ChatRequest):
    messages = [SystemMessage(content=(
        "You are UTKARSH AI, an elite productivity assistant. "
        "Be concise, structured, and insightful."
    ))]
    for m in req.history[-10:]:   # last 10 turns for context
        cls = HumanMessage if m["role"] == "user" else AIMessage
        messages.append(cls(content=m["content"]))
    messages.append(HumanMessage(content=req.message))

    def stream():
        for chunk in llm.stream(messages):
            yield chunk.content

    return StreamingResponse(stream(), media_type="text/plain")


# ═══════════════════════════════════════════════════════════════
#  2.  PPT GENERATION
# ═══════════════════════════════════════════════════════════════
DARK_BG  = RGBColor(0x0D, 0x0D, 0x1A)
ACCENT   = RGBColor(0x7C, 0x3A, 0xED)   # violet
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
SUBTEXT  = RGBColor(0xA0, 0xA0, 0xC0)

def _set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_title_slide(prs, topic: str):
    layout = prs.slide_layouts[6]          # blank
    slide  = prs.slides.add_slide(layout)
    _set_bg(slide, DARK_BG)
    W, H = prs.slide_width, prs.slide_height

    # Accent bar
    bar = slide.shapes.add_shape(1, 0, H - Inches(0.08), W, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Title
    txb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.4))
    tf  = txb.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = topic.upper()
    run.font.size  = Pt(44); run.font.bold = True
    run.font.color.rgb = WHITE

    # Subtitle
    txb2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(6), Inches(0.6))
    tf2  = txb2.text_frame
    p2   = tf2.paragraphs[0]
    run2 = p2.add_run(); run2.text = "Generated by UTKARSH AI"
    run2.font.size = Pt(18); run2.font.color.rgb = SUBTEXT

def _add_content_slide(prs, title: str, bullets: list[str], slide_num: int):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _set_bg(slide, DARK_BG)
    W, H = prs.slide_width, prs.slide_height

    # Slide number dot
    circ = slide.shapes.add_shape(9, W - Inches(0.9), Inches(0.15), Inches(0.5), Inches(0.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    txc = circ.text_frame.paragraphs[0]
    txc.alignment = PP_ALIGN.CENTER
    rc = txc.add_run(); rc.text = str(slide_num)
    rc.font.size = Pt(11); rc.font.bold = True; rc.font.color.rgb = WHITE

    # Title
    txb = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.4), Inches(0.9))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = title
    run.font.size = Pt(28); run.font.bold = True; run.font.color.rgb = WHITE

    # Divider
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.25), Inches(3), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Bullets
    txb2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.8), Inches(4.8))
    tf2  = txb2.text_frame; tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        para = tf2.add_paragraph() if i else tf2.paragraphs[0]
        para.space_before = Pt(6)
        run  = para.add_run()
        run.text = f"  ▸  {bullet}"
        run.font.size  = Pt(17)
        run.font.color.rgb = WHITE if i % 2 == 0 else SUBTEXT


@app.post("/generate-ppt")
async def generate_ppt(req: PPTRequest):
    # 1. Ask LLM for slide outline
    outline_prompt = f"""
    Create a {req.num_slides}-slide presentation outline on: "{req.topic}"
    Return ONLY valid JSON in this exact format, no markdown:
    {{
      "slides": [
        {{"title": "...", "bullets": ["...", "...", "...", "..."]}}
      ]
    }}
    Each slide must have 3-5 bullet points. Make it informative and professional.
    """
    resp = llm.invoke([HumanMessage(content=outline_prompt)])
    raw  = resp.content.strip()
    # Strip possible ```json fences
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"): raw = raw[4:]
    outline = json.loads(raw.strip())

    # 2. Build PPTX
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, req.topic)
    for i, s in enumerate(outline["slides"], start=1):
        _add_content_slide(prs, s["title"], s["bullets"], i)

    # Thank-you slide
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _set_bg(slide, DARK_BG)
    txb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = "Thank You"
    run.font.size = Pt(52); run.font.bold = True; run.font.color.rgb = ACCENT
    p2  = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2  = p2.add_run(); r2.text = "Generated by UTKARSH AI"
    r2.font.size = Pt(18); r2.font.color.rgb = SUBTEXT

    filename = f"utkarsh_{uuid.uuid4().hex[:8]}.pptx"
    path     = OUTPUT_DIR / filename
    prs.save(str(path))

    return FileResponse(
        path=str(path),
        filename=filename,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# ═══════════════════════════════════════════════════════════════
#  3.  DOCUMENT Q&A  (RAG with ChromaDB)
# ═══════════════════════════════════════════════════════════════
_vector_stores: dict[str, Chroma] = {}

@app.post("/upload-doc")
async def upload_doc(file: UploadFile = File(...)):
    ext  = Path(file.filename).suffix.lower()
    path = UPLOAD_DIR / f"{uuid.uuid4().hex}{ext}"
    path.write_bytes(await file.read())

    loader = PyPDFLoader(str(path)) if ext == ".pdf" else Docx2txtLoader(str(path))
    docs   = loader.load()

    splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=100)
    chunks   = splitter.split_documents(docs)

    session = uuid.uuid4().hex[:12]
    vs = Chroma.from_documents(
        chunks, embeddings,
        persist_directory=str(CHROMA_DIR / session)
    )
    _vector_stores[session] = vs

    return {"session_id": session, "chunks": len(chunks), "filename": file.filename}


@app.post("/doc-qa")
async def doc_qa(
    session_id: str = Form(...),
    question:   str = Form(...),
):
    vs = _vector_stores.get(session_id)
    if not vs:
        raise HTTPException(404, "Session not found. Upload a document first.")
    chain = RetrievalQA.from_chain_type(llm=llm, retriever=vs.as_retriever(search_kwargs={"k": 4}))
    answer = chain.invoke({"query": question})
    return {"answer": answer["result"]}


# ═══════════════════════════════════════════════════════════════
#  4.  IMAGE ANALYSIS  (Claude Vision)
# ═══════════════════════════════════════════════════════════════
@app.post("/analyze-image")
async def analyze_image(
    file:   UploadFile = File(...),
    prompt: str        = Form("Describe this image in detail and extract all visible text."),
):
    data      = await file.read()
    b64_image = base64.standard_b64encode(data).decode("utf-8")
    mime      = file.content_type or "image/jpeg"

    message = ant_client.messages.create(
        model="claude-opus-4-5",
        max_tokens=1024,
        messages=[{
            "role": "user",
            "content": [
                {"type": "image", "source": {"type": "base64", "media_type": mime, "data": b64_image}},
                {"type": "text", "text": prompt},
            ],
        }],
    )
    return {"analysis": message.content[0].text}


# ═══════════════════════════════════════════════════════════════
#  5.  SUMMARIZE DOCUMENT (quick, no RAG)
# ═══════════════════════════════════════════════════════════════
@app.post("/summarize-doc")
async def summarize_doc(file: UploadFile = File(...)):
    ext  = Path(file.filename).suffix.lower()
    path = UPLOAD_DIR / f"{uuid.uuid4().hex}{ext}"
    path.write_bytes(await file.read())

    loader = PyPDFLoader(str(path)) if ext == ".pdf" else Docx2txtLoader(str(path))
    docs   = loader.load()
    text   = " ".join(d.page_content for d in docs)[:12000]   # ~3k tokens

    resp = llm.invoke([
        SystemMessage(content="You are a precise document summarizer."),
        HumanMessage(content=f"Summarize this document concisely with key points:\n\n{text}"),
    ])
    return {"summary": resp.content, "filename": file.filename}


# ═══════════════════════════════════════════════════════════════
if __name__ == "__main__":
    import uvicorn
    uvicorn.run("utkarsh_ai_backend:app", host="0.0.0.0", port=8000, reload=True)
